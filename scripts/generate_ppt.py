import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    DARK_BG = RGBColor(11, 12, 14)       # #0b0c0e
    CARD_BG = RGBColor(20, 23, 29)       # #14171d
    AMBER = RGBColor(217, 119, 6)        # #d97706
    WHITE = RGBColor(240, 240, 240)
    GRAY = RGBColor(160, 165, 175)
    GREEN = RGBColor(16, 185, 129)
    RED = RGBColor(239, 68, 68)

    blank_layout = prs.slide_layouts[6]
    
    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()

    def add_header(slide, title_text, category_text="PHYSIOXAI // NANO TECHNOLOGY HACKATHON"):
        # Category label
        tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.4))
        p_cat = tb_cat.text_frame.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = AMBER
        p_cat.font.name = 'Courier New'
        
        # Main Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11), Inches(0.8))
        p_title = tb_title.text_frame.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        p_title.font.name = 'Arial'

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)
    
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0))
    tf = tb.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "PHYSIOXAI"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = AMBER
    p1.font.name = 'Arial'
    
    p2 = tf.add_paragraph()
    p2.text = "Physics-Grounded Explainable AI for Physical Signal Classification"
    p2.font.size = Pt(24)
    p2.font.color.rgb = WHITE
    p2.font.name = 'Arial'
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Nano Technology Hackathon Project | Vibration Sensor Analytics Workstation"
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY
    p3.font.name = 'Courier New'
    p3.space_before = Pt(30)

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement & Motivation
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_header(slide2, "1. The Core Hackathon Problem")
    
    # Left Card
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = RED
    tf1 = card1.text_frame
    tf1.margin_left = tf1.margin_top = Inches(0.3)
    p = tf1.paragraphs[0]
    p.text = "❌ The AI Shortcut Trap"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = RED
    
    points1 = [
        "Traditional AI models report 99% accuracy on signal datasets.",
        "However, models often learn spurious noise or background artifacts rather than true physical mechanisms.",
        "In safety-critical applications (aviation, nuclear, cleanrooms), unverified black-box predictions are unacceptable."
    ]
    for pt in points1:
        p = tf1.add_paragraph()
        p.text = "• " + pt
        p.font.size = Pt(14); p.font.color.rgb = WHITE; p.space_before = Pt(14)

    # Right Card
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = GREEN
    tf2 = card2.text_frame
    tf2.margin_left = tf2.margin_top = Inches(0.3)
    p = tf2.paragraphs[0]
    p.text = "✅ The PhysioXAI Requirement"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = GREEN
    
    points2 = [
        "Explain AI predictions through measurable physical features.",
        "Perform Controlled Physical Intervention directly on raw signal frequency components.",
        "Re-evaluate modified signals using the EXACT SAME trained model.",
        "Provide reproducible empirical proof of model sensitivity."
    ]
    for pt in points2:
        p = tf2.add_paragraph()
        p.text = "• " + pt
        p.font.size = Pt(14); p.font.color.rgb = WHITE; p.space_before = Pt(14)

    # -------------------------------------------------------------
    # SLIDE 3: System Architecture & Anti-Cheating Engine
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_header(slide3, "2. Anti-Cheating Pipeline & Architecture")
    
    card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = AMBER
    tf = card.text_frame
    tf.margin_left = tf.margin_top = Inches(0.4)
    
    p = tf.paragraphs[0]
    p.text = "Strict Anti-Cheating Execution Flow"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = AMBER
    
    steps = [
        "1. Raw Vibration Signal Input → Single-Sided Fast Fourier Transform (FFT).",
        "2. Physical Feature Extraction → RMS, Peak, Crest Factor, Spectral Centroid & 100-140Hz Band Energy.",
        "3. Controlled Frequency Intervention → Multiply target band energy in frequency domain via FFT.",
        "4. Signal Synthesis → Reconstruct modified time-domain signal using Inverse FFT (IFFT).",
        "5. Same ML Model Evaluation → Re-extract features & evaluate blindly via static Random Forest artifact.",
        "6. Empirical Delta Output → Measure exact physical feature change vs. AI prediction probability shift."
    ]
    for st in steps:
        p = tf.add_paragraph()
        p.text = st
        p.font.size = Pt(14); p.font.color.rgb = WHITE; p.space_before = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 4: Physical Feature Engineering
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_header(slide4, "3. Domain Feature Extraction & FFT Analysis")
    
    features = [
        ("Characteristic Band Energy (100–140Hz)", "Primary fault indicator for bearing pass frequency harmonics.", AMBER),
        ("RMS Amplitude (Root Mean Square)", "Total signal power measure across time domain waveform.", WHITE),
        ("Peak Amplitude & Crest Factor", "Peak-to-RMS ratio highlighting transient impact spikes.", WHITE),
        ("Spectral Centroid & Total Energy", "Center of spectral mass and total frequency power distribution.", WHITE),
        ("Dominant Frequency (Hz)", "Peak magnitude frequency indicating baseline shaft rotation.", WHITE)
    ]
    
    top = 1.8
    for title, desc, color in features:
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.6), Inches(0.85))
        box.fill.solid(); box.fill.fore_color.rgb = CARD_BG; box.line.color.rgb = color
        tf = box.text_frame
        tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12); p2.font.color.rgb = GRAY

        top += 1.0

    # -------------------------------------------------------------
    # SLIDE 5: Workstation Dashboard Overview (Screenshot)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_header(slide5, "4. Workstation Interface: Baseline Signal Analysis")
    
    img_path1 = "docs/screenshots/slide_baseline_analysis.png"
    if os.path.exists(img_path1):
        slide5.shapes.add_picture(img_path1, Inches(0.8), Inches(1.8), width=Inches(11.6))

    # -------------------------------------------------------------
    # SLIDE 6: Controlled Physical Intervention (Screenshot)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)
    add_header(slide6, "5. Controlled Physical Intervention Engine")
    
    img_path2 = "docs/screenshots/slide_intervention_result.png"
    if os.path.exists(img_path2):
        slide6.shapes.add_picture(img_path2, Inches(0.8), Inches(1.8), width=Inches(11.6))

    # -------------------------------------------------------------
    # SLIDE 7: Scientific Methodology & Integrity Rules
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7)
    add_header(slide7, "6. Scientific Methodology & Integrity Rules")
    
    card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    card.fill.solid(); card.fill.fore_color.rgb = CARD_BG; card.line.color.rgb = AMBER
    tf = card.text_frame
    tf.margin_left = tf.margin_top = Inches(0.4)
    
    p = tf.paragraphs[0]
    p.text = "Strict Scientific Wording Guidelines"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = AMBER
    
    rules = [
        ("✅ APPROVED SCIENTIFIC CLAIMS", GREEN, [
            "\"The system demonstrates sensitivity to a measurable physical feature through controlled intervention.\"",
            "\"Controlled manipulation of the characteristic frequency produces a measurable change in classifier output.\"",
            "\"The experiment provides evidence that the classifier responds to the manipulated physical feature.\""
        ]),
        ("🚫 PROHIBITED SCIENTIFIC CLAIMS", RED, [
            "NEVER claim: \"Our AI understands physics.\"",
            "NEVER claim: \"The AI has discovered causality.\"",
            "NEVER claim: \"The model is 100% guaranteed to use the physical feature in all scenarios.\""
        ])
    ]
    
    for section_title, color, items in rules:
        p = tf.add_paragraph()
        p.text = section_title
        p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = color; p.space_before = Pt(14)
        for item in items:
            p = tf.add_paragraph()
            p.text = "  • " + item
            p.font.size = Pt(13); p.font.color.rgb = WHITE; p.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 8: Technology Stack & Implementation
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8)
    add_header(slide8, "7. Technology Stack & Engineering Architecture")
    
    cols = [
        ("Frontend Workstation", [
            "React 18 & TypeScript",
            "Vite & Tailwind CSS",
            "Recharts Interactive Charts",
            "Axios API Integration",
            "Industrial Dark Palette"
        ], AMBER),
        ("Backend & Signal Math", [
            "Python 3.10 & FastAPI",
            "Pydantic v2 Validation",
            "NumPy & SciPy FFT Math",
            "Pandas Data Wrangling",
            "Uvicorn Async Gateway"
        ], WHITE),
        ("ML & Security Layer", [
            "Random Forest Classifier",
            "Joblib Server Model Load",
            "Safe Memory Stream CSV",
            "Strict CORS & CSP Headers",
            "Render Cloud Deployment"
        ], GREEN)
    ]
    
    left = 0.8
    for col_title, items, color in cols:
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_BG; card.line.color.rgb = color
        tf = card.text_frame; tf.margin_left = tf.margin_top = Inches(0.3)
        p = tf.paragraphs[0]
        p.text = col_title
        p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = color
        for it in items:
            p = tf.add_paragraph()
            p.text = "• " + it
            p.font.size = Pt(13); p.font.color.rgb = WHITE; p.space_before = Pt(12)
        left += 4.0

    # -------------------------------------------------------------
    # SLIDE 9: Security Architecture & Validation
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9)
    add_header(slide9, "8. Security Architecture & Input Hardening")
    
    card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    card.fill.solid(); card.fill.fore_color.rgb = CARD_BG; card.line.color.rgb = AMBER
    tf = card.text_frame; tf.margin_left = tf.margin_top = Inches(0.4)
    
    p = tf.paragraphs[0]
    p.text = "Security Protection Mechanisms"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = AMBER
    
    sec_points = [
        "1. Input Validation & Bounds: Pydantic schemas enforce bounds on signal length (max 50k points), NaN checks, & sampling rate limits.",
        "2. Safe CSV Parsing: Multi-format memory parser prevents arbitrary code execution and rejects non-CSV / binary upload attacks.",
        "3. Model Artifact Security: Ban on user model uploads; loads only repository-controlled, verified Random Forest model artifacts.",
        "4. Path Traversal & Command Protection: Absolute isolation of server paths and zero execution of user-controlled shell subprocesses.",
        "5. Headers & CORS Security: Custom middleware injecting X-Content-Type-Options, CSP, and restricted CORS origin rules."
    ]
    for sp in sec_points:
        p = tf.add_paragraph()
        p.text = sp
        p.font.size = Pt(13.5); p.font.color.rgb = WHITE; p.space_before = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 10: Real-World Applications & Impact
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10)
    add_header(slide10, "9. Real-World Applications & Nanotech Impact")
    
    apps = [
        ("Precision Nanomanufacturing", "Validating vibration sensor diagnostics in cleanroom nano-fabrication tools to prevent wafer damage.", AMBER),
        ("Predictive Machinery Maintenance", "Isolating bearing and gear defect frequencies in turbines and pumps before catastrophic failure.", WHITE),
        ("Safety-Critical Infrastructure", "Auditing AI diagnostic models in aviation and power plants to eliminate shortcut learning risk.", WHITE),
        ("Sensor Engineering Verification", "Providing sensor manufacturers with empirical verification of classifier sensitivity under noisy environments.", WHITE)
    ]
    
    top = 1.8
    for title, desc, color in apps:
        box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.6), Inches(1.1))
        box.fill.solid(); box.fill.fore_color.rgb = CARD_BG; box.line.color.rgb = color
        tf = box.text_frame; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13); p2.font.color.rgb = GRAY; p2.space_before = Pt(4)
        top += 1.25

    # -------------------------------------------------------------
    # SLIDE 11: Experimental Verification & Results
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide11)
    add_header(slide11, "10. Experimental Verification & Test Results")
    
    card = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    card.fill.solid(); card.fill.fore_color.rgb = CARD_BG; card.line.color.rgb = GREEN
    tf = card.text_frame; tf.margin_left = tf.margin_top = Inches(0.4)
    
    p = tf.paragraphs[0]
    p.text = "Verification & Validation Metrics"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = GREEN
    
    res = [
        "• Model Training Performance: Random Forest achieved 100% Accuracy, Precision, Recall, & F1-Score on test split.",
        "• Pytest Suite Pass Rate: 100% pass rate across signal FFT calculations, feature extraction, model inference, and intervention engine.",
        "• Anti-Cheating Test Verification: Confirmed empirical shift in recomputed signal features and model predictions.",
        "• Live Deployment Verification: Deployed on Render Static Site + Web Service with full CORS & API integration."
    ]
    for r in res:
        p = tf.add_paragraph()
        p.text = r
        p.font.size = Pt(14); p.font.color.rgb = WHITE; p.space_before = Pt(14)

    # -------------------------------------------------------------
    # SLIDE 12: Conclusion & Summary
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide12)
    add_header(slide12, "11. Conclusion & Summary")
    
    card = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.6), Inches(4.8))
    card.fill.solid(); card.fill.fore_color.rgb = CARD_BG; card.line.color.rgb = AMBER
    tf = card.text_frame; tf.margin_left = tf.margin_top = Inches(0.4)
    
    p = tf.paragraphs[0]
    p.text = "PhysioXAI Core Summary"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = AMBER
    
    summary_pts = [
        "1. Solved Hackathon Problem: Connects AI predictions to measurable physical features via Controlled Intervention.",
        "2. Proven Anti-Cheating Math: Frequency-domain signal modification (FFT/IFFT) evaluated blindly by the same ML model.",
        "3. Scientifically Responsible: Strict adherence to defensible claims without over-promising AI capabilities.",
        "4. Full-Stack Production Quality: High data-density React dashboard, secure FastAPI backend, 100% tested & deployed live."
    ]
    for s in summary_pts:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(15); p.font.color.rgb = WHITE; p.space_before = Pt(14)

    output_path = "docs/PhysioXAI_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
